import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
RED = RGBColor(0xE7, 0x4C, 0x3C)
DARK_RED = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xF3, 0x9C, 0x12)
GREEN = RGBColor(0x27, 0xAE, 0x60)
DARK = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFD, 0xF6, 0xEC)
LIGHT_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
CARD_BG = RGBColor(0xFF, 0xF5, 0xEE)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_rect(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ==================== Slide 1: Cover ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

# Decorative circles
add_circle(slide, -0.5, -0.5, 3, RGBColor(0xE7, 0x4C, 0x3C))
add_circle(slide, 11.5, 5.5, 2.5, RGBColor(0xF3, 0x9C, 0x12))
add_circle(slide, 10.8, -0.8, 1.8, RGBColor(0x27, 0xAE, 0x60))

# Title
add_textbox(slide, 1.5, 1.8, 10.3, 1.2, '🍜 南宁美食推荐网站', font_size=48, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.0, 10.3, 0.8, '舌尖上的绿城 · 美食指南', font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()

add_textbox(slide, 1.5, 4.2, 10.3, 0.6, 'AI 驱动的全栈项目实战', font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.0, 10.3, 0.5, '从零到上线 · 完整开发流程演示', font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 6.3, 10.3, 0.5, '2026年6月  |  liuclick.github.io/nanning-food', font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ==================== Slide 2: Project Overview ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Title bar
add_rect(slide, 0, 0, 13.333, 1.2, DARK_RED)
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '📋 项目概述', font_size=32, color=WHITE, bold=True)

add_textbox(slide, 0.8, 1.8, 11.7, 0.6, '项目背景', font_size=22, color=RED, bold=True)
add_textbox(slide, 0.8, 2.4, 11.7, 0.8, '为用户打造一个展示南宁特色美食的公开网站，收录12道经典美食，涵盖五大分类，可在线浏览、搜索、筛选，支持手机端访问。', font_size=16, color=DARK)

add_textbox(slide, 0.8, 3.4, 11.7, 0.6, '核心目标', font_size=22, color=RED, bold=True)

items = [
    '✅ 收录南宁12道经典美食，每道菜配详细文字介绍',
    '✅ 提供分类筛选、实时搜索、详情弹窗等交互功能',
    '✅ 3条美食路线 + 6大饮食文化板块，内容丰富',
    '✅ 响应式设计，适配手机、平板、电脑',
    '✅ 永久免费部署到公网，可随时访问分享',
]
for i, item in enumerate(items):
    add_textbox(slide, 1.2, 4.0 + i * 0.5, 10.5, 0.5, item, font_size=15, color=DARK)

# ==================== Slide 3: Tech Stack ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, 0, 0, 13.333, 1.2, DARK)
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '🛠️ 技术栈', font_size=32, color=WHITE, bold=True)

techs = [
    ('HTML5', '页面结构，语义化标签', '🟠', RGBColor(0xE6, 0x5C, 0x00)),
    ('CSS3', '渐变背景、动画、响应式布局、Grid/Flexbox', '🔵', RGBColor(0x29, 0x6E, 0xD9)),
    ('JavaScript', '数据渲染、搜索筛选、弹窗交互、事件处理', '🟡', RGBColor(0xF7, 0xDF, 0x1E)),
    ('Git', '版本控制，代码管理', '⚫', RGBColor(0x2C, 0x3E, 0x50)),
    ('GitHub', '代码托管，在线仓库', '⬛', RGBColor(0x24, 0x23, 0x2E)),
    ('GitHub Pages', '免费静态网站托管，永久公网访问', '🟢', RGBColor(0x27, 0xAE, 0x60)),
]
for i, (name, desc, icon, clr) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.8 + row * 2.6
    
    card = add_rect(slide, x, y, 3.7, 2.2, RGBColor(0xFA, 0xFA, 0xFA))
    add_rect(slide, x, y, 3.7, 0.06, clr)
    add_textbox(slide, x + 0.3, y + 0.3, 3.1, 0.5, f'{icon} {name}', font_size=20, color=DARK, bold=True)
    add_textbox(slide, x + 0.3, y + 0.9, 3.1, 1.0, desc, font_size=14, color=LIGHT_GRAY)

# ==================== Slide 4: Features ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, 0, 0, 13.333, 1.2, RED)
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '✨ 核心功能一览', font_size=32, color=WHITE, bold=True)

features = [
    ('🔍', '智能搜索', '支持美食名称、描述、分类\n的模糊实时搜索'),
    ('🏷️', '分类筛选', '5大分类标签，一键过滤\n米粉/硬菜/小吃/甜品/酸嘢'),
    ('🪟', '详情弹窗', '点击卡片弹出详情窗口\n含介绍、推荐店铺、评分'),
    ('🗺️', '美食路线', '3条精心规划的一日游\n美食探索路线推荐'),
    ('📖', '美食文化', '6大板块介绍南宁饮食\n文化特色与饮食传统'),
    ('📱', '响应式设计', '手机/平板/电脑完美适配\n移动端汉堡菜单导航'),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.8 + row * 2.6
    
    card = add_rect(slide, x, y, 3.7, 2.2, CARD_BG)
    add_textbox(slide, x + 0.3, y + 0.2, 3.1, 0.5, icon, font_size=28, color=DARK)
    add_textbox(slide, x + 0.3, y + 0.7, 3.1, 0.4, title, font_size=18, color=DARK, bold=True)
    add_textbox(slide, x + 0.3, y + 1.2, 3.1, 0.8, desc, font_size=13, color=LIGHT_GRAY)

# ==================== Slide 5: Food Content ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, 0, 0, 13.333, 1.2, GREEN)
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '🍽️ 美食内容矩阵', font_size=32, color=WHITE, bold=True)

foods_data = [
    ('🍜 米粉类', '老友粉、生榨米粉、卷筒粉', '南宁人的灵魂食物，一日三餐离不开'),
    ('🥘 硬菜大餐', '柠檬鸭、横县鱼生、南宁烧烤', '宴客硬菜，南宁饮食文化的精髓'),
    ('🥟 小吃点心', '粉饺、五色糯米饭', '街头巷尾的经典味道'),
    ('🥒 酸嘢腌渍', '酸嘢', '英雄难过美人关，美人难过酸嘢摊'),
    ('🍹 甜品饮品', '凉茶、龟苓膏、八宝饭', '清热解暑，药食同源的智慧'),
]
for i, (cat, items_list, desc) in enumerate(foods_data):
    x = 0.5 + i * 2.5
    card = add_rect(slide, x, 1.8, 2.2, 3.6, CARD_BG)
    add_textbox(slide, x + 0.15, 2.0, 1.9, 0.5, cat, font_size=16, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 2.6, 1.9, 0.8, items_list, font_size=13, color=DARK, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.3, 3.5, 1.6, 0.03, GOLD)
    add_textbox(slide, x + 0.15, 3.7, 1.9, 1.4, desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 5.8, 11.7, 0.5, '每道美食包含：名称、Emoji、分类标签、辣度等级、价格区间、推荐店铺、评分、简介、详细介绍', font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Stats
stats = [('12', '道美食'), ('5', '大分类'), ('3', '条路线'), ('6', '个文化板块')]
for i, (num, label) in enumerate(stats):
    x = 2.5 + i * 2.5
    add_textbox(slide, x, 6.4, 1.8, 0.5, num, font_size=36, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 6.9, 1.8, 0.4, label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ==================== Slide 6: Design Highlights ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, 0, 0, 13.333, 1.2, GOLD)
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '🎨 设计亮点', font_size=32, color=WHITE, bold=True)

designs = [
    ('温暖米色调', '整体配色以温暖的米色、奶油色为主，搭配红、金、绿点缀，营造亲切温暖的美食氛围。'),
    ('渐变卡片', '每道美食卡片使用独特的双色渐变背景，搭配大号Emoji和点阵纹理，即使没有图片也视觉出彩。'),
    ('悬停动画', '卡片悬停时Emoji放大旋转、卡片上浮阴影加深，弹窗平滑过渡，细节处处精致。'),
    ('固定导航栏', '毛玻璃效果的固定导航栏，滚动时自动收缩，点击平滑跳转，体验流畅。'),
    ('移动端适配', '汉堡菜单、卡片自适应、触摸优化，手机端浏览体验同样出色。'),
]
for i, (title, desc) in enumerate(designs):
    y = 1.6 + i * 1.1
    add_rect(slide, 0.8, y, 0.08, 0.8, RED)
    add_textbox(slide, 1.2, y, 3.0, 0.4, title, font_size=18, color=DARK, bold=True)
    add_textbox(slide, 1.2, y + 0.4, 10.5, 0.5, desc, font_size=14, color=LIGHT_GRAY)

# ==================== Slide 7: Deployment ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, 0, 0, 13.333, 1.2, DARK)
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '🚀 部署流程', font_size=32, color=WHITE, bold=True)

steps = [
    ('1', '本地开发', '编写 HTML/CSS/JS\n单文件完成所有功能'),
    ('2', 'Git 初始化', '安装 Git，初始化仓库\ngit init & git commit'),
    ('2', 'Git 初始化', '安装 Git，初始化仓库\ngit init & git commit'),
    ('3', 'GitHub 仓库', '创建 GitHub 仓库\n推送代码到远程'),
    ('4', 'GitHub Pages', '开启 GitHub Pages\n自动部署上线'),
    ('5', '永久访问', '获得永久公网 URL\n随时分享给任何人'),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.6 + i * 2.5
    # Number circle
    add_circle(slide, x + 0.7, 1.8, 0.65, RED)
    add_textbox(slide, x + 0.7, 1.85, 0.65, 0.55, num, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow (except last)
    if i < 4:
        add_textbox(slide, x + 1.6, 1.85, 0.6, 0.55, '→', font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    # Content
    add_textbox(slide, x + 0.15, 2.7, 2.2, 0.4, title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 3.2, 2.2, 0.9, desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Result URL
add_rect(slide, 0.8, 4.5, 11.7, 1.5, CARD_BG)
add_textbox(slide, 1.2, 4.7, 10.9, 0.5, '🎯 最终成果', font_size=20, color=RED, bold=True)
add_textbox(slide, 1.2, 5.2, 10.9, 0.5, '🔗 https://liuclick.github.io/nanning-food/', font_size=22, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ==================== Slide 8: AI Collaboration ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, 0, 0, 13.333, 1.2, RGBColor(0x8E, 0x44, 0xAD))
add_textbox(slide, 0.8, 0.2, 11.7, 0.8, '🤖 AI 协作开发过程', font_size=32, color=WHITE, bold=True)

ai_steps = [
    ('需求分析', '用户提出"制作南宁美食推荐网址"，AI分析需求并设计完整方案'),
    ('代码生成', 'AI一次性生成完整HTML页面，包含CSS样式、JS逻辑、美食数据'),
    ('迭代优化', '根据用户反馈添加分类标签、搜索功能、响应式设计、回到顶部等'),
    ('图片方案', '从外部图片API → 纯CSS渐变+Emoji方案，解决图片不匹配问题'),
    ('部署上线', '安装Git → 创建GitHub仓库 → 配置GitHub Pages → 获得永久URL'),
    ('PPT生成', '使用python-pptx库自动生成项目演示PPT，总结全过程'),
]
for i, (title, desc) in enumerate(ai_steps):
    y = 1.5 + i * 0.95
    add_rect(slide, 0.8, y, 0.5, 0.5, RED)
    add_textbox(slide, 0.8, y + 0.05, 0.5, 0.4, str(i+1), font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.6, y, 2.2, 0.4, title, font_size=17, color=DARK, bold=True)
    add_textbox(slide, 4.0, y, 8.5, 0.7, desc, font_size=14, color=LIGHT_GRAY)

# ==================== Slide 9: Summary ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_circle(slide, -0.5, -0.5, 3, RGBColor(0xE7, 0x4C, 0x3C))
add_circle(slide, 11.5, 5.5, 2.5, RGBColor(0xF3, 0x9C, 0x12))

add_textbox(slide, 1.5, 0.8, 10.3, 0.8, '🎉 项目总结', font_size=40, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, 5.5, 1.7, 2.3, 0.04, RED)

summary_items = [
    '📄 单文件 HTML 实现完整美食推荐网站',
    '🎨 纯 CSS 渐变 + Emoji 打造精美视觉',
    '📱 响应式设计，三端完美适配',
    '🔍 搜索 + 分类 + 弹窗，交互丰富',
    '🚀 GitHub Pages 永久免费托管',
    '🤖 AI 全程辅助，从零到上线',
    '🌐 公网可访问：liuclick.github.io/nanning-food',
]
for i, item in enumerate(summary_items):
    add_textbox(slide, 2.5, 2.2 + i * 0.6, 8.3, 0.5, item, font_size=18, color=DARK, alignment=PP_ALIGN.LEFT)

add_textbox(slide, 1.5, 6.5, 10.3, 0.5, '感谢观看！🍜 欢迎访问网站体验', font_size=20, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# Save
output_path = r'c:\Users\86159\Desktop\新建文件夹\南宁美食推荐项目演示.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')