from pptx import Presentation
from pptx.util import Pt

path = r'c:\Users\86159\Desktop\新建文件夹\AI考研自习室-项目演示.pptx'
prs = Presentation(path)

for slide_idx in [0, 1, 2, 3, 28]:
    slide = prs.slides[slide_idx]
    print(f'=== 第{slide_idx+1}页 ===')
    for shape in slide.shapes:
        left_in = round(shape.left / 914400, 2)
        top_in = round(shape.top / 914400, 2)
        w_in = round(shape.width / 914400, 2)
        h_in = round(shape.height / 914400, 2)
        
        fill_color = 'none'
        try:
            if shape.fill.type is not None:
                fill_color = str(shape.fill.fore_color.rgb)
        except:
            fill_color = 'theme'
        
        texts = []
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    font_info = ''
                    if p.runs:
                        r = p.runs[0]
                        try:
                            font_info = f'[sz={r.font.size},b={r.font.bold}]'
                        except:
                            pass
                    texts.append(f'{t[:60]}{font_info}')
        
        print(f'  {shape.name}: pos=({left_in},{top_in}) {w_in}x{h_in}, fill={fill_color}')
        for t in texts:
            print(f'    "{t}"')
    print()