from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import os

# Template colors
DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
LIGHT_BLUE = RGBColor(0xBE, 0xE3, 0xF8)
ACCENT_BLUE = RGBColor(0x31, 0x82, 0xCE)
DEEP_BLUE = RGBColor(0x2C, 0x52, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x2D, 0x37, 0x48)
GRAY_TEXT = RGBColor(0x71, 0x80, 0x96)

PROJECT_NAME = '南宁美食推荐'
PROJECT_YEAR = '2026'

# Load template
template_path = r'c:\Users\86159\Desktop\新建文件夹\AI考研自习室-项目演示.pptx'
prs = Presentation(template_path)

# Delete all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

layout = prs.slide_layouts[0]  # DEFAULT layout

def add_slide():
    return prs.slides.add_slide(layout)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline_textbox(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """lines is list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        text, font_size, color, bold, alignment = line_info
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
    return tf

def add_header_bar(slide, page_num, section_title=''):
    # Top bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 10, 0.55, DARK_BLUE)
    # Circle in header
    add_shape(slide, MSO_SHAPE.OVAL, 0.3, 0.1, 0.35, 0.35, LIGHT_BLUE)
    # Project name
    add_textbox(slide, 0.75, 0.1, 1.5, 0.35, PROJECT_NAME, 11, WHITE, True)
    # Section tabs
    tabs = ['项目概述', '美食内容', '功能特色', '设计亮点', '部署上线', '项目总结']
    for i, tab in enumerate(tabs):
        x = 2.0 + i * 1.2
        if i == section_title:
            add_textbox(slide, x, 0.05, 1.0, 0.35, tab, 10, WHITE, True)
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.15, 0.42, 0.7, 0.04, ACCENT_BLUE)
        else:
            add_textbox(slide, x, 0.05, 1.0, 0.35, tab, 10, RGBColor(0x90, 0xAE, 0xD4), False)
    # Footer
    add_textbox(slide, 0.5, 5.0, 3.0, 0.3, f'{PROJECT_NAME} · {PROJECT_YEAR}', 8, GRAY_TEXT)
    # Page number
    pg_shape = add_shape(slide, MSO_SHAPE.OVAL, 9.3, 5.1, 0.4, 0.4, DEEP_BLUE)
    add_textbox(slide, 9.3, 5.1, 0.4, 0.4, str(page_num), 11, WHITE, True, PP_ALIGN.CENTER)

def add_card(slide, left, top, width, height, emoji='', title='', desc='', title_size=14, desc_size=11):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, LIGHT_BLUE)
    if emoji:
        add_textbox(slide, left + width/2 - 0.3, top + 0.15, 0.6, 0.45, emoji, 22, BLACK, False, PP_ALIGN.CENTER)
    if title:
        add_textbox(slide, left + 0.15, top + 0.65, width - 0.3, 0.3, title, title_size, DARK_TEXT, True, PP_ALIGN.CENTER)
    if desc:
        add_textbox(slide, left + 0.15, top + 0.95, width - 0.3, height - 1.1, desc, desc_size, GRAY_TEXT, False, PP_ALIGN.CENTER)

# ==================== Slide 1: Cover ====================
slide = add_slide()
# Decorative circles
add_shape(slide, MSO_SHAPE.OVAL, 4.0, 0.8, 2.0, 2.0, DEEP_BLUE)
add_shape(slide, MSO_SHAPE.OVAL, 7.5, 3.5, 0.8, 0.8, ACCENT_BLUE)
add_shape(slide, MSO_SHAPE.OVAL, 0.3, 4.0, 0.5, 0.5, LIGHT_BLUE)
# Title
add_textbox(slide, 0, 1.5, 10, 1.2, '🍜 南宁美食推荐', 48, DARK_TEXT, True, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.6, 10, 0.6, '舌尖上的绿城 · 美食指南', 20, GRAY_TEXT, False, PP_ALIGN.CENTER)
# Divider
add_shape(slide, MSO_SHAPE.RECTANGLE, 4.0, 3.3, 2.0, 0.03, ACCENT_BLUE)
add_textbox(slide, 0, 3.5, 10, 0.5, 'AI辅助全栈开发 · 从零到上线', 16, DEEP_BLUE, True, PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.2, 10, 0.35, 'liuclick.github.io/nanning-food', 14, GRAY_TEXT, False, PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.7, 10, 0.35, '2026年6月', 12, GRAY_TEXT, False, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 5.15, 3.0, 0.25, f'{PROJECT_NAME} · {PROJECT_YEAR}', 8, GRAY_TEXT)

# ==================== Slide 2: TOC ====================
slide = add_slide()
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 10, 0.55, DARK_BLUE)
add_textbox(slide, 0.5, 0.1, 3.0, 0.35, f'{PROJECT_NAME} · {PROJECT_YEAR}', 11, WHITE, True)
add_textbox(slide, 0, 0.8, 10, 0.4, '目  录', 36, DARK_TEXT, True, PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, 4.3, 1.25, 1.4, 0.03, ACCENT_BLUE)
add_textbox(slide, 0, 1.5, 10, 0.4, 'C O N T E N T S', 14, GRAY_TEXT, False, PP_ALIGN.CENTER)

toc_items = [
    ('01', '项目概述', '项目背景与核心目标'),
    ('02', '美食内容', '12道美食 · 5大分类'),
    ('03', '功能特色', '搜索/筛选/弹窗/路线'),
    ('04', '设计亮点', '视觉风格与交互体验'),
    ('05', '部署上线', 'GitHub Pages 永久托管'),
    ('06', '项目总结', 'AI协作过程回顾'),
]
for i, (num, title, desc) in enumerate(toc_items):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 3.1
    y = 2.2 + row * 1.5
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.8, 1.2, LIGHT_BLUE)
    add_textbox(slide, x + 0.15, y + 0.1, 0.5, 0.4, num, 22, ACCENT_BLUE, True)
    add_textbox(slide, x + 0.65, y + 0.15, 2.0, 0.3, title, 14, DARK_TEXT, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.65, y + 0.5, 1.2, 0.02, ACCENT_BLUE)
    add_textbox(slide, x + 0.65, y + 0.6, 2.0, 0.4, desc, 10, GRAY_TEXT)

add_textbox(slide, 0.5, 5.15, 3.0, 0.25, f'{PROJECT_NAME} · {PROJECT_YEAR}', 8, GRAY_TEXT)

# ==================== Slide 3: Project Overview ====================
slide = add_slide()
add_header_bar(slide, 3, 0)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 项目概述 · 背景与目标', 18, DARK_TEXT, True)

# Pain points
cards_data = [
    ('😰', '信息碎片化', '南宁美食信息分散在各平台，游客难以快速找到靠谱的本地美食推荐'),
    ('📉', '缺乏整合', '没有一个平台将南宁美食按分类、辣度、价格、位置进行系统化整理'),
    ('💡', '核心洞察', '打造一个简洁美观的南宁美食指南网站，一站式解决游客和本地人的觅食需求'),
]
for i, (emoji, title, desc) in enumerate(cards_data):
    x = 0.6 + i * 3.1
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.3, 2.8, 2.2, LIGHT_BLUE)
    add_textbox(slide, x + 1.1, 1.5, 0.6, 0.5, emoji, 20, BLACK, False, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 2.05, 2.5, 0.35, title, 14, DARK_TEXT, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 2.45, 2.5, 0.9, desc, 10, GRAY_TEXT)

# Bottom insight
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 3.75, 6.5, 0.85, LIGHT_BLUE)
add_textbox(slide, 0.8, 3.85, 6.1, 0.6, '南宁拥有丰富的美食文化，老友粉、柠檬鸭、酸嘢等特色美食享誉全国。但市面上缺少一个专门面向南宁美食的、美观实用的在线指南。', 11, DARK_TEXT)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.24, 3.6, 2.26, 1.12, WHITE)
add_multiline_textbox(slide, 7.34, 3.6, 2.06, 1.12, [
    ('💡 核心目标', 11, DARK_TEXT, True, PP_ALIGN.LEFT),
    ('收录12道美食', 11, ACCENT_BLUE, True, PP_ALIGN.LEFT),
    ('永久免费在线访问', 11, ACCENT_BLUE, True, PP_ALIGN.LEFT),
])

# ==================== Slide 4: Tech Stack ====================
slide = add_slide()
add_header_bar(slide, 4, 0)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 项目概述 · 技术栈', 18, DARK_TEXT, True)

techs = [
    ('🟠', 'HTML5', '页面结构\n语义化标签'),
    ('🔵', 'CSS3', '渐变/动画/Grid\n响应式布局'),
    ('🟡', 'JavaScript', '数据渲染\n交互逻辑'),
    ('⚫', 'Git', '版本控制\n代码管理'),
    ('⬛', 'GitHub', '代码托管\n在线仓库'),
    ('🟢', 'GitHub Pages', '免费托管\n永久公网访问'),
]
for i, (emoji, name, desc) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 3.1
    y = 1.3 + row * 2.0
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.8, 1.6, LIGHT_BLUE)
    add_textbox(slide, x + 0.2, y + 0.15, 2.4, 0.35, f'{emoji} {name}', 14, DARK_TEXT, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.2, y + 0.55, 1.0, 0.02, ACCENT_BLUE)
    add_textbox(slide, x + 0.2, y + 0.65, 2.4, 0.8, desc, 10, GRAY_TEXT)

# ==================== Slide 5: Food Content ====================
slide = add_slide()
add_header_bar(slide, 5, 1)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 美食内容 · 12道经典南宁美食', 18, DARK_TEXT, True)

food_cats = [
    ('🍜', '米粉类', '老友粉\n生榨米粉\n卷筒粉', '南宁灵魂食物'),
    ('🥘', '硬菜大餐', '柠檬鸭\n横县鱼生\n南宁烧烤', '宴客必点硬菜'),
    ('🥟', '小吃点心', '粉饺\n五色糯米饭', '街头经典味道'),
    ('🥒', '酸嘢腌渍', '酸嘢', '酸甜脆爽开胃'),
    ('🍹', '甜品饮品', '凉茶 龟苓膏\n八宝饭', '清热解暑养生'),
]
for i, (emoji, name, items, tag) in enumerate(food_cats):
    x = 0.6 + i * 1.9
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.3, 1.6, 3.0, LIGHT_BLUE)
    add_textbox(slide, x + 0.05, 1.45, 1.5, 0.35, f'{emoji} {name}', 13, DARK_TEXT, True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.3, 1.85, 1.0, 0.02, ACCENT_BLUE)
    add_textbox(slide, x + 0.05, 2.0, 1.5, 1.2, items, 10, GRAY_TEXT, False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.2, 3.5, 1.2, 0.35, DEEP_BLUE)
    add_textbox(slide, x + 0.2, 3.5, 1.2, 0.35, tag, 8, WHITE, True, PP_ALIGN.CENTER)

# Stats
stats = [('12', '道美食'), ('5', '大分类'), ('3', '条路线'), ('6', '文化板块'), ('⭐⭐⭐⭐⭐', '高评分')]
for i, (num, label) in enumerate(stats):
    x = 0.8 + i * 1.9
    add_textbox(slide, x, 4.55, 1.5, 0.35, num, 24, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, 4.9, 1.5, 0.25, label, 9, GRAY_TEXT, False, PP_ALIGN.CENTER)

# ==================== Slide 6: Feature Highlights ====================
slide = add_slide()
add_header_bar(slide, 6, 2)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 功能特色 · 交互体验', 18, DARK_TEXT, True)

features = [
    ('🔍', '智能搜索', '支持美食名称、描述、分类的模糊实时搜索'),
    ('🏷️', '分类筛选', '5大分类标签，一键过滤，快速定位'),
    ('🪟', '详情弹窗', '点击卡片弹出详情，含介绍、店铺、评分'),
    ('🗺️', '美食路线', '3条精心规划的一日游美食探索路线'),
    ('📖', '美食文化', '6大板块介绍南宁饮食文化特色'),
    ('📱', '响应式设计', '手机/平板/电脑完美适配，汉堡菜单'),
]
for i, (emoji, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 3.1
    y = 1.3 + row * 2.0
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.8, 1.6, LIGHT_BLUE)
    add_textbox(slide, x + 0.2, y + 0.15, 2.4, 0.35, f'{emoji} {title}', 14, DARK_TEXT, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.2, y + 0.55, 1.0, 0.02, ACCENT_BLUE)
    add_textbox(slide, x + 0.2, y + 0.65, 2.4, 0.8, desc, 10, GRAY_TEXT)

# ==================== Slide 7: Design Highlights ====================
slide = add_slide()
add_header_bar(slide, 7, 3)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 设计亮点 · 视觉与交互', 18, DARK_TEXT, True)

designs = [
    ('温暖米色调', '整体配色以温暖米色、奶油色为主，红、金、绿点缀，营造亲切的美食氛围。'),
    ('渐变卡片', '每道美食使用独特双色渐变背景+大号Emoji+点阵纹理，视觉出彩。'),
    ('悬停动画', '卡片悬停Emoji放大旋转、卡片上浮、阴影加深，细节精致。'),
    ('固定导航栏', '毛玻璃效果固定导航，滚动自动收缩，点击平滑跳转。'),
    ('移动端适配', '汉堡菜单、卡片自适应、触摸优化，手机端体验出色。'),
]
for i, (title, desc) in enumerate(designs):
    y = 1.3 + i * 0.8
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 8.8, 0.65, LIGHT_BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, y, 0.05, 0.65, ACCENT_BLUE)
    add_textbox(slide, 0.85, y + 0.05, 2.0, 0.3, title, 14, DARK_TEXT, True)
    add_textbox(slide, 0.85, y + 0.35, 8.3, 0.25, desc, 10, GRAY_TEXT)

# ==================== Slide 8: Deployment ====================
slide = add_slide()
add_header_bar(slide, 8, 4)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 部署上线 · GitHub Pages 永久托管', 18, DARK_TEXT, True)

steps = [
    ('1', '本地开发', '编写HTML/CSS/JS\n单文件完成所有功能'),
    ('2', 'Git初始化', '安装Git → git init\n→ git commit'),
    ('3', 'GitHub仓库', '创建仓库 → 推送代码\n→ 生成Token认证'),
    ('4', 'GitHub Pages', '开启Pages功能\n自动构建部署'),
    ('5', '永久访问', '获得永久URL\n随时分享给任何人'),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.5 + i * 1.9
    # Number circle
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.55, 1.3, 0.5, 0.5, DEEP_BLUE)
    add_textbox(slide, x + 0.55, 1.3, 0.5, 0.5, num, 16, WHITE, True, PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        add_textbox(slide, x + 1.25, 1.3, 0.4, 0.5, '→', 18, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    # Card
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.15, 2.0, 1.6, 1.5, LIGHT_BLUE)
    add_textbox(slide, x + 0.25, 2.1, 1.4, 0.3, title, 12, DARK_TEXT, True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.45, 2.45, 0.7, 0.02, ACCENT_BLUE)
    add_textbox(slide, x + 0.25, 2.55, 1.4, 0.8, desc, 9, GRAY_TEXT, False, PP_ALIGN.CENTER)

# Result URL
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 3.8, 8.8, 0.8, LIGHT_BLUE)
add_textbox(slide, 0.8, 3.85, 1.0, 0.3, '🎯 最终成果', 12, DARK_TEXT, True)
add_textbox(slide, 0.8, 4.15, 8.4, 0.35, '🔗 https://liuclick.github.io/nanning-food/', 16, DEEP_BLUE, True, PP_ALIGN.CENTER)

# ==================== Slide 9: AI Collaboration ====================
slide = add_slide()
add_header_bar(slide, 9, 5)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 项目总结 · AI协作开发过程', 18, DARK_TEXT, True)

ai_steps = [
    ('需求分析', '用户提出"制作南宁美食推荐网址"，AI分析需求并设计完整方案'),
    ('代码生成', 'AI一次性生成完整HTML页面，包含CSS样式、JS逻辑、12道美食数据'),
    ('迭代优化', '根据反馈添加分类标签、搜索功能、响应式设计、回到顶部、图片方案等'),
    ('部署上线', '安装Git → 创建GitHub仓库 → 配置Pages → 获得永久URL'),
    ('PPT生成', '使用python-pptx库自动生成项目演示PPT，总结全过程'),
]
for i, (title, desc) in enumerate(ai_steps):
    y = 1.3 + i * 0.75
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 8.8, 0.6, LIGHT_BLUE)
    add_shape(slide, MSO_SHAPE.OVAL, 0.75, y + 0.1, 0.4, 0.4, DEEP_BLUE)
    add_textbox(slide, 0.75, y + 0.1, 0.4, 0.4, str(i+1), 14, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1.35, y + 0.05, 1.8, 0.25, title, 13, DARK_TEXT, True)
    add_textbox(slide, 1.35, y + 0.32, 7.8, 0.25, desc, 10, GRAY_TEXT)

# ==================== Slide 10: Summary ====================
slide = add_slide()
add_header_bar(slide, 10, 5)
add_textbox(slide, 0.5, 0.75, 8.8, 0.4, '▶ 项目总结 · 成果回顾', 18, DARK_TEXT, True)

summary_items = [
    '📄 单文件 HTML 实现完整美食推荐网站',
    '🎨 纯 CSS 渐变 + Emoji 打造精美视觉设计',
    '📱 响应式设计，手机/平板/电脑三端完美适配',
    '🔍 搜索 + 分类筛选 + 详情弹窗，交互丰富',
    '🚀 GitHub Pages 永久免费托管，公网可访问',
    '🤖 AI 全程辅助开发，从零到上线完整流程',
]
for i, item in enumerate(summary_items):
    y = 1.3 + i * 0.55
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 8.8, 0.45, LIGHT_BLUE)
    add_textbox(slide, 0.8, y + 0.05, 8.4, 0.35, item, 13, DARK_TEXT, False)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 4.65, 8.8, 0.5, DEEP_BLUE)
add_textbox(slide, 0.6, 4.65, 8.8, 0.5, '🌐  公网访问地址：https://liuclick.github.io/nanning-food/', 14, WHITE, True, PP_ALIGN.CENTER)

# ==================== Slide 11: Thank You ====================
slide = add_slide()
add_shape(slide, MSO_SHAPE.OVAL, 4.0, 0.8, 2.0, 2.0, DEEP_BLUE)
add_shape(slide, MSO_SHAPE.OVAL, 7.5, 3.5, 0.8, 0.8, ACCENT_BLUE)
add_textbox(slide, 0, 1.8, 10, 1.2, '感谢聆听', 48, DARK_TEXT, True, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.9, 10, 0.6, 'T H A N K   Y O U', 18, GRAY_TEXT, False, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 3.5, 9.0, 0.5, '"用美食感受南宁的城市温度"', 14, GRAY_TEXT, False, PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.8, 10, 0.35, f'项目: {PROJECT_NAME}  |  访问: liuclick.github.io/nanning-food  |  周期: 2026年6月', 10, GRAY_TEXT, False, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 5.15, 3.0, 0.25, f'{PROJECT_NAME} · {PROJECT_YEAR}', 8, GRAY_TEXT)

# Save
output_path = r'c:\Users\86159\Desktop\新建文件夹\南宁美食推荐项目演示.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')